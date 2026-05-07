from __future__ import annotations

import base64
import logging
import mimetypes
import queue
import time
import traceback
import uuid
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from threading import Event, Thread
from types import SimpleNamespace
from typing import Any, Iterator
from xml.etree import ElementTree as ET

from sqlalchemy import func, select
from sqlalchemy.orm import Session

from app.core.config import get_settings
from app.db.database import session_scope
from app.db.models import ChatAttachment, ChatMessageRecord, ChatSession
from app.schemas.aniu import (
    ChatAttachmentRead,
    ChatMessageRead,
    ChatSessionRead,
    ChatStreamRequest,
)
from app.skills.providers import build_skill_context
from app.services.llm_service import LLMStreamCancelled, llm_service


logger = logging.getLogger(__name__)

MAX_UPLOAD_BYTES = 10 * 1024 * 1024  # 10 MB
MAX_ATTACHMENTS_PER_MESSAGE = 12
UPLOAD_URL_PREFIX = "/api/aniu/chat/uploads"
DEFAULT_SESSION_TITLE = "新对话"
ATTACHMENT_TEXT_LIMIT = 12_000
ATTACHMENT_TOTAL_TEXT_LIMIT = 24_000
IMAGE_DATA_URL_MAX_BYTES = 6 * 1024 * 1024
ALLOWED_APPLICATION_MIME_TYPES = {
    "application/json",
    "application/ld+json",
    "application/xml",
    "application/yaml",
    "application/x-yaml",
    "application/x-ndjson",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
ALLOWED_FILE_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".tsv",
    ".json",
    ".jsonl",
    ".xml",
    ".yaml",
    ".yml",
    ".log",
    ".ini",
    ".cfg",
    ".conf",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".vue",
    ".java",
    ".go",
    ".rs",
    ".c",
    ".cc",
    ".cpp",
    ".h",
    ".hpp",
    ".sql",
    ".sh",
    ".ps1",
    ".bat",
    ".docx",
    ".xlsx",
    ".pptx",
}


def _assume_utc(value: datetime | None) -> datetime | None:
    if value is None:
        return None
    if value.tzinfo is None:
        return value.replace(tzinfo=timezone.utc)
    return value.astimezone(timezone.utc)


def _uploads_root(session_id: int | None = None) -> Path:
    settings = get_settings()
    root = settings.sqlite_db_path.parent / "chat_uploads"
    if session_id is not None:
        root = root / str(session_id)
    root.mkdir(parents=True, exist_ok=True)
    return root


def _attachment_to_read(attachment: ChatAttachment) -> ChatAttachmentRead:
    return ChatAttachmentRead(
        id=attachment.id,
        filename=attachment.filename,
        mime_type=attachment.mime_type,
        size=attachment.size,
        url=f"{UPLOAD_URL_PREFIX}/{attachment.id}",
    )


def _attachment_dict(attachment: ChatAttachment) -> dict[str, Any]:
    return {
        "id": attachment.id,
        "filename": attachment.filename,
        "mime_type": attachment.mime_type,
        "size": attachment.size,
        "url": f"{UPLOAD_URL_PREFIX}/{attachment.id}",
        "storage_path": attachment.storage_path,
    }


def _session_to_read(session: ChatSession, message_count: int = 0) -> ChatSessionRead:
    return ChatSessionRead(
        id=session.id,
        title=session.title,
        kind=str(session.kind or "user"),
        slug=session.slug,
        created_at=_assume_utc(session.created_at),
        updated_at=_assume_utc(session.updated_at),
        last_message_at=_assume_utc(session.last_message_at),
        message_count=message_count,
    )


def _message_to_read(record: ChatMessageRecord) -> ChatMessageRead:
    attachments_payload = record.attachments or []
    attachments = [
        ChatAttachmentRead(**item)
        for item in attachments_payload
        if isinstance(item, dict)
    ]
    return ChatMessageRead(
        id=record.id,
        role=record.role,
        content=record.content,
        tool_calls=record.tool_calls,
        attachments=attachments or None,
        created_at=_assume_utc(record.created_at),
    )


def _normalize_attachment_type(filename: str, mime_type: str) -> tuple[str, str]:
    safe_filename = Path(filename).name or "upload.bin"
    suffix = Path(safe_filename).suffix.lower()
    normalized_mime = str(mime_type or "").strip().lower()
    guessed_mime, _ = mimetypes.guess_type(safe_filename)

    if not normalized_mime or normalized_mime == "application/octet-stream":
        normalized_mime = str(guessed_mime or "application/octet-stream").lower()

    if normalized_mime.startswith("image/"):
        return safe_filename, normalized_mime
    if normalized_mime.startswith("text/"):
        return safe_filename, normalized_mime
    if normalized_mime in ALLOWED_APPLICATION_MIME_TYPES:
        return safe_filename, normalized_mime
    if suffix in ALLOWED_FILE_EXTENSIONS:
        fallback_mime = normalized_mime
        if fallback_mime == "application/octet-stream" and guessed_mime:
            fallback_mime = guessed_mime.lower()
        if fallback_mime == "application/octet-stream":
            fallback_mime = "text/plain"
        return safe_filename, fallback_mime

    raise ValueError("仅支持图片、文本以及 docx/xlsx/pptx 等现代办公文档附件。")


def _attachment_prompt_text(item: dict[str, Any]) -> str | None:
    filename = str(item.get("filename") or "attachment").strip() or "attachment"
    mime_type = str(item.get("mime_type") or "").strip().lower()
    url = str(item.get("url") or "").strip()
    if mime_type.startswith("image/"):
        if url:
            return f"![{filename}]({url})"
        return f"[用户上传了图片：{filename}]"
    return f"[用户上传了文件：{filename}]"


def _trim_text(text: str, limit: int) -> tuple[str, bool]:
    normalized = str(text or "").replace("\r\n", "\n").strip()
    if len(normalized) <= limit:
        return normalized, False
    return normalized[:limit].rstrip(), True


def _read_text_with_fallback(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _extract_docx_text(path: Path) -> str:
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with zipfile.ZipFile(path) as archive:
        with archive.open("word/document.xml") as document_xml:
            root = ET.parse(document_xml).getroot()

    paragraphs: list[str] = []
    for paragraph in root.findall(".//w:body/w:p", namespace):
        texts = [node.text or "" for node in paragraph.findall(".//w:t", namespace)]
        merged = "".join(texts).strip()
        if merged:
            paragraphs.append(merged)
    return "\n".join(paragraphs)


def _extract_xlsx_text(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    sheets: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for values in sheet.iter_rows(values_only=True):
            cells = [str(value).strip() for value in values if value not in {None, ""}]
            if cells:
                rows.append("\t".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    workbook.close()
    return "\n\n".join(sheets)


def _extract_pptx_text(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
        if lines:
            slides.append(f"[Slide {index}]\n" + "\n".join(lines))
    return "\n\n".join(slides)


def _extract_attachment_text(path: Path, mime_type: str) -> str:
    normalized = str(mime_type or "").strip().lower()
    suffix = path.suffix.lower()

    if normalized.startswith("text/") or suffix in {
        ".txt",
        ".md",
        ".markdown",
        ".csv",
        ".tsv",
        ".json",
        ".jsonl",
        ".xml",
        ".yaml",
        ".yml",
        ".log",
        ".ini",
        ".cfg",
        ".conf",
        ".py",
        ".js",
        ".ts",
        ".tsx",
        ".jsx",
        ".vue",
        ".java",
        ".go",
        ".rs",
        ".c",
        ".cc",
        ".cpp",
        ".h",
        ".hpp",
        ".sql",
        ".sh",
        ".ps1",
        ".bat",
    }:
        return _read_text_with_fallback(path)
    if (
        normalized
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or suffix == ".docx"
    ):
        return _extract_docx_text(path)
    if (
        normalized
        == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or suffix == ".xlsx"
    ):
        return _extract_xlsx_text(path)
    if (
        normalized
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or suffix == ".pptx"
    ):
        return _extract_pptx_text(path)
    raise ValueError("暂不支持提取该文件类型的正文。")


class ChatSessionService:
    def _build_attachment_content_parts(
        self,
        attachments_payload: list[dict[str, Any]] | None,
    ) -> list[dict[str, Any]]:
        if not attachments_payload:
            return []

        content_parts: list[dict[str, Any]] = []
        remaining_text_budget = ATTACHMENT_TOTAL_TEXT_LIMIT
        for item in attachments_payload:
            if not isinstance(item, dict):
                continue
            filename = str(item.get("filename") or "attachment").strip() or "attachment"
            mime_type = str(item.get("mime_type") or "").strip().lower()
            storage_path = str(item.get("storage_path") or "").strip()
            if not storage_path:
                prompt_text = _attachment_prompt_text(item)
                if prompt_text:
                    content_parts.append({"type": "text", "text": prompt_text})
                continue

            path = Path(storage_path)
            if not path.is_file():
                content_parts.append(
                    {
                        "type": "text",
                        "text": f"[附件 {filename}] 文件缺失，无法读取内容。",
                    }
                )
                continue

            if mime_type.startswith("image/"):
                raw = path.read_bytes()
                if len(raw) > IMAGE_DATA_URL_MAX_BYTES:
                    content_parts.append(
                        {
                            "type": "text",
                            "text": (
                                f"[图片附件 {filename}] 文件过大，已跳过图片直传。"
                            ),
                        }
                    )
                    continue
                data_url = (
                    f"data:{mime_type or 'application/octet-stream'};base64,"
                    f"{base64.b64encode(raw).decode('ascii')}"
                )
                content_parts.append(
                    {
                        "type": "image_url",
                        "image_url": {"url": data_url},
                    }
                )
                continue

            if remaining_text_budget <= 0:
                content_parts.append(
                    {
                        "type": "text",
                        "text": (
                            f"[附件 {filename}] 由于附件文本总量超出上限，后续内容已省略。"
                        ),
                    }
                )
                continue

            try:
                extracted = _extract_attachment_text(path, mime_type)
                if not extracted.strip():
                    content_parts.append(
                        {
                            "type": "text",
                            "text": f"[附件 {filename}] 未提取到可用正文。",
                        }
                    )
                    continue
                attachment_limit = min(ATTACHMENT_TEXT_LIMIT, remaining_text_budget)
                trimmed, truncated = _trim_text(extracted, attachment_limit)
                remaining_text_budget = max(0, remaining_text_budget - len(trimmed))
                note = "\n\n[内容已截断]" if truncated else ""
                content_parts.append(
                    {
                        "type": "text",
                        "text": (
                            f"[附件 {filename}] 以下为提取文本：\n{trimmed}{note}"
                        ),
                    }
                )
            except Exception as exc:  # noqa: BLE001
                logger.warning(
                    "attachment text extraction failed: filename=%s mime_type=%s error=%s",
                    filename,
                    mime_type,
                    exc,
                )
                content_parts.append(
                    {
                        "type": "text",
                        "text": f"[附件 {filename}] 无法提取正文：{exc}",
                    }
                )
        return content_parts

    def _build_user_message_content(
        self,
        *,
        content: str,
        attachments_payload: list[dict[str, Any]] | None,
    ) -> str | list[dict[str, Any]]:
        text = str(content or "").strip()
        attachment_parts = self._build_attachment_content_parts(attachments_payload)
        if not attachment_parts:
            return text

        content_parts: list[dict[str, Any]] = []
        if text:
            content_parts.append({"type": "text", "text": text})
        content_parts.extend(attachment_parts)
        return content_parts

    def list_sessions(self, db: Session) -> list[ChatSessionRead]:
        count_sub = (
            select(
                ChatMessageRecord.session_id,
                func.count(ChatMessageRecord.id).label("count"),
            )
            .group_by(ChatMessageRecord.session_id)
            .subquery()
        )

        rows = db.execute(
            select(ChatSession, func.coalesce(count_sub.c.count, 0))
            .where(ChatSession.kind == "user")
            .outerjoin(count_sub, count_sub.c.session_id == ChatSession.id)
            .order_by(
                ChatSession.last_message_at.desc().nullslast(),
                ChatSession.updated_at.desc(),
            )
        ).all()

        return [_session_to_read(session, int(count)) for session, count in rows]

    def create_session(
        self, db: Session, *, title: str | None = None
    ) -> ChatSessionRead:
        session = ChatSession(
            title=(title or DEFAULT_SESSION_TITLE).strip() or DEFAULT_SESSION_TITLE,
            kind="user",
        )
        db.add(session)
        db.flush()
        return _session_to_read(session, 0)

    def rename_session(
        self, db: Session, session_id: int, *, title: str
    ) -> ChatSessionRead:
        session = db.get(ChatSession, session_id)
        if session is None or str(session.kind or "user") != "user":
            raise LookupError("会话不存在。")
        session.title = title.strip() or session.title
        db.flush()
        count = db.execute(
            select(func.count(ChatMessageRecord.id)).where(
                ChatMessageRecord.session_id == session_id
            )
        ).scalar_one()
        return _session_to_read(session, int(count))

    def delete_session(self, db: Session, session_id: int) -> None:
        session = db.get(ChatSession, session_id)
        if session is None or str(session.kind or "user") != "user":
            raise LookupError("会话不存在。")
        db.delete(session)

    def list_messages(
        self,
        db: Session,
        session_id: int,
        *,
        limit: int = 50,
        before_id: int | None = None,
    ) -> tuple[ChatSessionRead, list[ChatMessageRead], int | None, bool]:
        session = db.get(ChatSession, session_id)
        if session is None or str(session.kind or "user") != "user":
            raise LookupError("会话不存在。")

        page_size = max(1, int(limit))
        total_count = db.execute(
            select(func.count(ChatMessageRecord.id)).where(
                ChatMessageRecord.session_id == session_id
            )
        ).scalar_one()

        stmt = (
            select(ChatMessageRecord)
            .where(ChatMessageRecord.session_id == session_id)
            .order_by(ChatMessageRecord.id.desc())
        )
        if before_id is not None:
            stmt = stmt.where(ChatMessageRecord.id < before_id)

        records = (
            db.execute(
                stmt.limit(page_size + 1)
            )
            .scalars()
            .all()
        )
        has_more = len(records) > page_size
        if has_more:
            records = records[:page_size]
        records.reverse()
        next_before_id = records[0].id if has_more and records else None
        return (
            _session_to_read(session, int(total_count)),
            [_message_to_read(record) for record in records],
            next_before_id,
            has_more,
        )

    def save_attachment(
        self,
        db: Session,
        *,
        filename: str,
        mime_type: str,
        data: bytes,
        session_id: int | None = None,
    ) -> ChatAttachmentRead:
        if len(data) > MAX_UPLOAD_BYTES:
            raise ValueError(
                f"文件过大，最大允许 {MAX_UPLOAD_BYTES // (1024 * 1024)}MB。"
            )

        safe_filename, normalized_mime_type = _normalize_attachment_type(
            filename, mime_type
        )
        suffix = Path(safe_filename).suffix
        stored_name = f"{uuid.uuid4().hex}{suffix}"
        storage_path = _uploads_root(session_id) / stored_name
        storage_path.write_bytes(data)

        attachment = ChatAttachment(
            filename=safe_filename,
            mime_type=normalized_mime_type,
            size=len(data),
            storage_path=str(storage_path),
        )
        db.add(attachment)
        db.flush()
        return _attachment_to_read(attachment)

    def get_attachment_file(
        self, db: Session, attachment_id: int
    ) -> tuple[Path, str, str]:
        attachment = db.get(ChatAttachment, attachment_id)
        if attachment is None:
            raise LookupError("附件不存在。")
        path = Path(attachment.storage_path)
        if not path.is_file():
            raise LookupError("附件文件丢失。")
        return path, attachment.mime_type, attachment.filename

    def _resolve_attachments(
        self, db: Session, attachment_ids: list[int]
    ) -> list[ChatAttachment]:
        if not attachment_ids:
            return []
        unique_ids = list(dict.fromkeys(int(x) for x in attachment_ids))
        records = (
            db.execute(
                select(ChatAttachment).where(ChatAttachment.id.in_(unique_ids))
            )
            .scalars()
            .all()
        )
        by_id = {record.id: record for record in records}
        resolved: list[ChatAttachment] = []
        for attachment_id in unique_ids:
            record = by_id.get(attachment_id)
            if record is None:
                raise LookupError(f"附件 {attachment_id} 不存在。")
            resolved.append(record)
        return resolved

    def _build_history_messages(
        self, records: list[ChatMessageRecord]
    ) -> list[dict[str, Any]]:
        history: list[dict[str, Any]] = []
        for record in records:
            if record.role not in {"user", "assistant", "system"}:
                continue

            if record.role == "user":
                history.append(
                    {
                        "role": record.role,
                        "content": self._build_user_message_content(
                            content=record.content,
                            attachments_payload=record.attachments,
                        ),
                    }
                )
                continue

            text = str(record.content or "").strip()
            if record.attachments:
                content_parts = [text] if text else []
                for item in record.attachments:
                    if not isinstance(item, dict):
                        continue
                    prompt_text = _attachment_prompt_text(item)
                    if prompt_text:
                        content_parts.append(prompt_text)
                text = "\n".join(part for part in content_parts if part).strip()

            entry: dict[str, Any] = {"role": record.role, "content": text}
            if record.role == "assistant":
                meta = record.meta_payload if isinstance(record.meta_payload, dict) else {}
                reasoning_content = meta.get("reasoning_content")
                if isinstance(reasoning_content, str) and reasoning_content.strip():
                    entry["reasoning_content"] = reasoning_content
            history.append(entry)
        return history

    def _derive_title(self, content: str) -> str:
        text = (content or "").strip().splitlines()
        first_line = text[0] if text else ""
        cleaned = first_line.strip() or "含附件消息"
        return cleaned[:30]

    def _build_failed_assistant_content(
        self,
        *,
        final_content: str,
        failed_message: str,
    ) -> str:
        content = str(final_content or "").strip()
        reason = str(failed_message or "聊天失败").strip() or "聊天失败"
        if content:
            return f"{content}\n\n执行失败：{reason}"
        return f"执行失败：{reason}"

    def _build_interrupted_assistant_content(self, final_content: str) -> str:
        content = str(final_content or "").strip()
        notice = "执行中断：客户端连接已断开。"
        if content:
            return f"{content}\n\n{notice}"
        return notice

    def stream_chat(self, payload: ChatStreamRequest) -> Iterator[dict[str, Any]]:
        if len(payload.attachment_ids) > MAX_ATTACHMENTS_PER_MESSAGE:
            raise ValueError(
                f"单条消息最多允许 {MAX_ATTACHMENTS_PER_MESSAGE} 个附件。"
            )

        with session_scope() as db:
            from app.services.aniu_service import aniu_service

            settings = aniu_service.get_or_create_settings(db)
            if not settings.llm_base_url or not settings.llm_api_key:
                raise RuntimeError("未配置大模型接口，无法执行 AI 聊天。")

            session = db.get(ChatSession, payload.session_id)
            if session is None or str(session.kind or "user") != "user":
                raise LookupError("会话不存在。")

            attachments = self._resolve_attachments(db, payload.attachment_ids)
            attachment_payload = [_attachment_dict(item) for item in attachments]
            normalized_content = str(payload.content or "").strip()

            user_record = ChatMessageRecord(
                session_id=session.id,
                role="user",
                content=normalized_content,
                attachments=attachment_payload or None,
            )
            db.add(user_record)
            db.flush()

            session.last_message_at = datetime.now(timezone.utc)
            if not session.title or session.title in {DEFAULT_SESSION_TITLE, "新会话"}:
                session.title = self._derive_title(normalized_content)

            history_records = (
                db.execute(
                    select(ChatMessageRecord)
                    .where(ChatMessageRecord.session_id == session.id)
                    .order_by(ChatMessageRecord.id.asc())
                )
                .scalars()
                .all()
            )
            history_messages = self._build_history_messages(history_records)
            settings_snapshot = SimpleNamespace(
                mx_api_key=settings.mx_api_key,
                system_prompt=settings.system_prompt,
                llm_model=settings.llm_model,
                llm_base_url=str(settings.llm_base_url),
                llm_api_key=str(settings.llm_api_key),
            )
            session_id = session.id

        event_queue: queue.Queue[dict[str, Any] | None] = queue.Queue()
        captured_tool_calls: list[dict[str, Any]] = []
        final_message_holder: dict[str, Any] = {}
        cancel_event = Event()

        def _emit(event_type: str, **data: Any) -> None:
            event_queue.put({"type": event_type, "ts": time.time(), **data})

        def _worker() -> None:
            try:
                result = llm_service.chat_result(
                    model=settings_snapshot.llm_model,
                    base_url=settings_snapshot.llm_base_url,
                    api_key=settings_snapshot.llm_api_key,
                    system_prompt=settings_snapshot.system_prompt,
                    messages=history_messages,
                    timeout_seconds=180,
                    tool_context=build_skill_context(
                        run_type="chat",
                        app_settings=settings_snapshot,
                    ),
                    emit=_emit,
                    cancel_event=cancel_event,
                )
                final_message_holder["final_message"] = result.get("final_message")
                content = str(result.get("final_answer") or "").strip()
                _emit("completed", message=content)
            except LLMStreamCancelled:
                logger.info("chat_session stream worker cancelled: session_id=%s", session_id)
            except Exception as exc:  # noqa: BLE001
                logger.exception("chat_session stream worker failed")
                _emit(
                    "failed",
                    message=str(exc),
                    traceback=traceback.format_exc(limit=4),
                )
            finally:
                event_queue.put(None)

        worker = Thread(
            target=_worker, daemon=True, name=f"aniu-chat-session-{session_id}"
        )
        worker.start()

        final_content = ""
        failed_message: str | None = None
        client_disconnected = False
        terminal_event_seen = False

        try:
            while True:
                try:
                    event = event_queue.get(timeout=15.0)
                except queue.Empty:
                    yield {"type": "heartbeat", "ts": time.time()}
                    continue
                if event is None:
                    break

                event_type = event.get("type")
                if event_type == "tool_call":
                    tool_name = str(event.get("tool_name") or "")
                    tool_call_id = str(event.get("tool_call_id") or "")
                    status = str(event.get("status") or "running")
                    existing = next(
                        (
                            item
                            for item in captured_tool_calls
                            if item["status"] == "running"
                            and (
                                (tool_call_id and item.get("tool_call_id") == tool_call_id)
                                or (
                                    not tool_call_id
                                    and item["tool_name"] == tool_name
                                )
                            )
                        ),
                        None,
                    )
                    if status == "running" and existing is None:
                        captured_tool_calls.append(
                            {
                                "tool_name": tool_name,
                                "tool_call_id": tool_call_id or None,
                                "status": "running",
                                "arguments": event.get("arguments"),
                                "started_at": event.get("ts"),
                            }
                        )
                    elif status == "done":
                        if existing is not None:
                            existing["status"] = "done"
                            existing["ok"] = event.get("ok")
                            existing["summary"] = event.get("summary")
                            existing["finished_at"] = event.get("ts")
                        else:
                            captured_tool_calls.append(
                                {
                                    "tool_name": tool_name,
                                    "tool_call_id": tool_call_id or None,
                                    "status": "done",
                                    "ok": event.get("ok"),
                                    "summary": event.get("summary"),
                                    "started_at": event.get("ts"),
                                    "finished_at": event.get("ts"),
                                }
                            )
                elif event_type == "final_delta":
                    delta = str(event.get("delta") or "")
                    if delta:
                        final_content += delta
                elif event_type in {"final_finished", "completed", "llm_final", "llm_message"}:
                    payload_content = str(
                        event.get("content") or event.get("message") or ""
                    )
                    if payload_content:
                        final_content = payload_content
                elif event_type == "failed":
                    failed_message = str(event.get("message") or "聊天失败")

                if event_type in {"completed", "failed"}:
                    terminal_event_seen = True

                yield event

                if terminal_event_seen:
                    while True:
                        trailing = event_queue.get()
                        if trailing is None:
                            break
                    break
        except GeneratorExit:
            client_disconnected = True
            raise
        finally:
            cancel_event.set()
            worker.join(timeout=1.0)

            with session_scope() as db:
                session = db.get(ChatSession, session_id)
                if session is not None:
                    assistant_content = ""
                    if failed_message is not None:
                        assistant_content = self._build_failed_assistant_content(
                            final_content=final_content,
                            failed_message=failed_message,
                        )
                    elif client_disconnected and not terminal_event_seen:
                        assistant_content = self._build_interrupted_assistant_content(
                            final_content
                        )
                    else:
                        assistant_content = str(final_content or "").strip()

                    meta_payload: dict[str, Any] | None = None
                    final_message = final_message_holder.get("final_message")
                    if isinstance(final_message, dict):
                        reasoning_value = final_message.get("reasoning_content") or final_message.get(
                            "reasoning"
                        )
                        if isinstance(reasoning_value, str) and reasoning_value.strip():
                            meta_payload = {"reasoning_content": reasoning_value}

                    if assistant_content or captured_tool_calls:
                        assistant_record = ChatMessageRecord(
                            session_id=session_id,
                            role="assistant",
                            content=assistant_content,
                            tool_calls=captured_tool_calls or None,
                            meta_payload=meta_payload,
                        )
                        db.add(assistant_record)
                        session.last_message_at = datetime.now(timezone.utc)


chat_session_service = ChatSessionService()
